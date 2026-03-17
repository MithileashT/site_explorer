"""Slack thread reader + local Ollama summarization service."""

from __future__ import annotations

import base64
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
from datetime import datetime, timezone
import io
import re
from typing import Dict, List, Tuple

import requests
from slack_sdk import WebClient
from slack_sdk.errors import SlackApiError

from core.config import resolve_slack_bot_token, settings
from core.logging import get_logger
from schemas.slack_investigation import (
    SlackLLMStatusResponse,
    SlackThreadAttachment,
    SlackThreadInvestigationRequest,
    SlackThreadInvestigationResponse,
    SlackThreadMessage,
)

logger = get_logger(__name__)

MAX_FILE_CHARS = 12_000
_URL_RE = re.compile(r"https?://\S+")
_MAX_ATTACHMENT_WORKERS = 4
_FILE_DOWNLOAD_TIMEOUT = 15  # seconds per file download


@dataclass
class ParsedSlackThreadRef:
    workspace: str
    channel_id: str
    thread_ts: str


def _p_timestamp_to_ts(value: str) -> str:
    if not value.isdigit() or len(value) < 7:
        raise ValueError("Invalid Slack message timestamp format.")
    return f"{value[:-6]}.{value[-6:]}"


def parse_slack_thread_url(url: str) -> ParsedSlackThreadRef:
    """Parse Slack permalink: .../archives/<CHANNEL>/p<TIMESTAMP>."""
    match = re.search(r"https://([^.]+)\.slack\.com/archives/([A-Z0-9]+)/p(\d+)", url)
    if not match:
        raise ValueError("Slack thread URL is invalid. Expected .../archives/<CHANNEL>/p<TIMESTAMP>.")
    workspace, channel_id, p_ts = match.groups()
    return ParsedSlackThreadRef(
        workspace=workspace,
        channel_id=channel_id,
        thread_ts=_p_timestamp_to_ts(p_ts),
    )


def _extract_log_blocks(text: str) -> Tuple[str, List[str]]:
    blocks: List[str] = []
    triple = re.findall(r"```(.*?)```", text or "", re.DOTALL)
    blocks.extend(b.strip() for b in triple if b.strip())
    clean = re.sub(r"```.*?```", "[log block]", text or "", flags=re.DOTALL)

    single = re.findall(r"`([^`]{40,})`", clean)
    blocks.extend(b.strip() for b in single if b.strip())
    clean = re.sub(r"`[^`]{40,}`", "[log snippet]", clean)
    return clean.strip(), blocks


def _as_bullets(text: str) -> List[str]:
    out: List[str] = []
    for line in text.splitlines():
        cleaned = re.sub(r"^[-*\d.\s]+", "", line).strip()
        if cleaned:
            out.append(cleaned)
    return out


def _split_markdown_sections(markdown_text: str) -> Dict[str, str]:
    sections: Dict[str, List[str]] = {}
    current = ""
    for line in markdown_text.splitlines():
        stripped = line.strip()
        # Match heading levels # through ####
        heading_match = re.match(r"^#{1,4}\s+(.+?)\s*$", stripped)
        if not heading_match:
            # Also match bold-only lines as headings (e.g. **The Issue**)
            heading_match = re.match(r"^\*\*(.+?)\*\*\s*$", stripped)
        if heading_match:
            # Normalize: lowercase, strip trailing colons/punctuation
            current = re.sub(r"[:\-]+$", "", heading_match.group(1).strip()).strip().lower()
            sections.setdefault(current, [])
            continue
        if current:
            sections[current].append(line)
    return {k: "\n".join(v).strip() for k, v in sections.items()}


def _find_section(sections: Dict[str, str], *candidates: str) -> str:
    """Fuzzy section lookup: try exact match, then substring containment."""
    for key in candidates:
        if key in sections:
            return sections[key].strip()
    # Fallback: check if any section key contains a candidate word
    for key in candidates:
        for sk, sv in sections.items():
            if key in sk or sk in key:
                return sv.strip()
    return ""


class SlackInvestigationService:
    def __init__(self, _llm_service=None) -> None:
        self._llm_service = _llm_service
        self.client: WebClient | None = None
        self._client_token = ""
        self._user_cache: Dict[str, str] = {}
        self._models_cache: List[str] | None = None
        self.ollama_host = settings.ollama_host.rstrip("/")
        self.text_model = settings.ollama_text_model
        # Reusable HTTP session for Slack file downloads (connection pooling)
        self._http_session: requests.Session | None = None

    def _slack_token(self) -> str:
        token = resolve_slack_bot_token()
        if token:
            return token
        raise RuntimeError(
            "SLACK_BOT_TOKEN is not configured on the backend. "
            "Set SLACK_BOT_TOKEN (or SLACK_TOKEN) in backend/.env. "
            "If using Docker, recreate backend: docker compose up -d --force-recreate backend"
        )

    def _require_client(self) -> WebClient:
        token = self._slack_token()
        if self.client is None or self._client_token != token:
            self.client = WebClient(token=token)
            self._client_token = token
        return self.client

    def _slack_headers(self) -> Dict[str, str]:
        return {"Authorization": f"Bearer {self._slack_token()}"}

    def _ollama_ping(self) -> bool:
        try:
            resp = requests.get(f"{self.ollama_host}/api/tags", timeout=3)
            return resp.status_code == 200
        except Exception:
            return False

    def _ollama_models(self) -> List[str]:
        if self._models_cache is not None:
            return self._models_cache
        try:
            resp = requests.get(f"{self.ollama_host}/api/tags", timeout=5)
            resp.raise_for_status()
            self._models_cache = [m.get("name", "") for m in resp.json().get("models", [])]
            return self._models_cache
        except Exception:
            return []

    def llm_status(self) -> SlackLLMStatusResponse:
        self._models_cache = None  # always fetch fresh for status check

        # Gather providers from LLMService if available
        providers: List[Dict] = []
        active_provider: Dict | None = None
        if self._llm_service and hasattr(self._llm_service, "available_providers"):
            providers = self._llm_service.available_providers()
            active_provider = self._llm_service.active_provider

        if not self._ollama_ping():
            # Ollama offline — still show OpenAI providers if available
            openai_providers = [p for p in providers if p.get("type") == "openai"]
            status = "offline" if not openai_providers else "online"
            fix_msg = (
                f"Ollama is not running at {self.ollama_host}. Run: ollama serve"
                if not openai_providers else None
            )
            return SlackLLMStatusResponse(
                status=status,
                text_model=self.text_model,
                text_ready=bool(openai_providers),
                installed=[],
                fix=fix_msg,
                providers=providers,
                active_provider=active_provider,
            )

        installed = self._ollama_models()
        text_prefix = self.text_model.split(":", 1)[0]
        return SlackLLMStatusResponse(
            status="online",
            text_model=self.text_model,
            text_ready=any(text_prefix in name for name in installed),
            installed=installed,
            providers=providers,
            active_provider=active_provider,
        )

    def _resolve_user(self, user_id: str | None) -> str:
        if not user_id:
            return "unknown"
        if user_id in self._user_cache:
            return self._user_cache[user_id]
        if user_id.startswith("B"):
            self._user_cache[user_id] = f"bot:{user_id}"
            return self._user_cache[user_id]

        client = self._require_client()
        try:
            data = client.users_info(user=user_id)
            user = data.get("user", {})
            profile = user.get("profile", {})
            display = profile.get("display_name") or profile.get("real_name") or user.get("name") or user_id
        except SlackApiError:
            display = user_id
        self._user_cache[user_id] = display
        return display

    def _download_file(self, url: str) -> bytes:
        if self._http_session is None:
            self._http_session = requests.Session()
            self._http_session.headers.update(self._slack_headers())
        resp = self._http_session.get(url, timeout=_FILE_DOWNLOAD_TIMEOUT)
        if resp.status_code != 200:
            raise ValueError(f"HTTP {resp.status_code}")
        return resp.content

    def _proc_image(self, data: bytes, filename: str) -> SlackThreadAttachment:
        return SlackThreadAttachment(
            filename=filename,
            filetype="image",
            extracted=f"[Image: {filename}]",
            b64_image=base64.b64encode(data).decode(),
        )

    def _proc_pdf(self, data: bytes, filename: str) -> SlackThreadAttachment:
        try:
            import pdfplumber

            pages: List[str] = []
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                for idx, page in enumerate(pdf.pages[:20], 1):  # cap at 20 pages
                    text = (page.extract_text() or "").strip()
                    if text:
                        pages.append(f"[Page {idx}]\n{text}")
            content = "\n\n".join(pages) or "[No extractable text in PDF]"
            return SlackThreadAttachment(
                filename=filename,
                filetype="pdf",
                extracted=content[:MAX_FILE_CHARS],
            )
        except Exception as exc:
            logger.warning("PDF extraction failed for %s: %s", filename, exc)
            return SlackThreadAttachment(filename=filename, filetype="pdf", extracted=f"[PDF read error: {exc}]")

    def _proc_pptx(self, data: bytes, filename: str) -> SlackThreadAttachment:
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            slides: List[str] = []
            for idx, slide in enumerate(prs.slides, 1):
                lines: List[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        for para in shape.text_frame.paragraphs:
                            line = " ".join(run.text for run in para.runs).strip()
                            if line:
                                lines.append(line)
                if lines:
                    slides.append(f"[Slide {idx}]\n" + "\n".join(lines))
            content = "\n\n".join(slides) or "[No text in presentation]"
            return SlackThreadAttachment(
                filename=filename,
                filetype="pptx",
                extracted=content[:MAX_FILE_CHARS],
            )
        except Exception as exc:
            logger.warning("PPTX extraction failed for %s: %s", filename, exc)
            return SlackThreadAttachment(filename=filename, filetype="pptx", extracted=f"[PPTX read error: {exc}]")

    def _proc_text(self, data: bytes, filename: str, filetype: str = "text") -> SlackThreadAttachment:
        content = data.decode("utf-8", errors="replace")
        return SlackThreadAttachment(
            filename=filename,
            filetype=filetype,
            extracted=content[:MAX_FILE_CHARS],
        )

    def _process_attachment(self, slack_file: Dict) -> SlackThreadAttachment | None:
        url = slack_file.get("url_private_download") or slack_file.get("url_private")
        if not url:
            return None

        filename = slack_file.get("name", "unknown")
        mimetype = slack_file.get("mimetype", "")
        lower_name = filename.lower()

        try:
            data = self._download_file(url)
        except Exception as exc:
            logger.warning("Attachment download failed for %s: %s", filename, exc)
            return SlackThreadAttachment(
                filename=filename,
                filetype="unknown",
                extracted=f"[Download failed: {exc}]",
            )

        if mimetype.startswith("image/"):
            return self._proc_image(data, filename)
        if mimetype == "application/pdf" or lower_name.endswith(".pdf"):
            return self._proc_pdf(data, filename)
        if "presentation" in mimetype or lower_name.endswith((".ppt", ".pptx")):
            return self._proc_pptx(data, filename)
        if mimetype.startswith("text/") or lower_name.endswith((".txt", ".log", ".yaml", ".yml", ".json", ".xml", ".csv")):
            filetype = "log" if lower_name.endswith(".log") else "text"
            return self._proc_text(data, filename, filetype)

        return SlackThreadAttachment(
            filename=filename,
            filetype="unknown",
            extracted=f"[Unsupported file type: {mimetype}]",
        )

    def _ensure_in_channel(self, client: WebClient, channel_id: str) -> None:
        """Join a public channel if the bot is not already a member."""
        try:
            channel_info = client.conversations_info(channel=channel_id)
            is_member = (channel_info.get("channel") or {}).get("is_member", False)
            if is_member:
                return
        except SlackApiError:
            pass  # proceed to join attempt

        try:
            client.conversations_join(channel=channel_id)
            logger.info("Auto-joined channel %s to fetch thread.", channel_id)
        except SlackApiError as exc:
            join_err = exc.response.get("error", "unknown")
            if join_err == "method_not_supported_for_channel_type":
                raise ValueError(
                    f"The bot is not in channel '{channel_id}' (private channel). "
                    "Invite the bot manually: /invite @<bot_name>"
                ) from exc
            if join_err == "is_archived":
                raise ValueError(
                    f"Channel '{channel_id}' is archived and cannot be joined."
                ) from exc
            raise ValueError(
                f"Could not join channel '{channel_id}' ({join_err}). "
                "Invite the bot to the channel first."
            ) from exc

    def _fetch_thread_messages(
        self,
        ref: ParsedSlackThreadRef,
        include_bots: bool,
        max_messages: int,
    ) -> Tuple[List[SlackThreadMessage], List[SlackThreadAttachment]]:
        client = self._require_client()
        cursor = None
        _auto_joined = False

        # ── Phase 1: Fetch raw Slack messages (no attachment download yet) ────
        raw_items: List[Dict] = []  # list of raw Slack message dicts
        while len(raw_items) < max_messages:
            limit = min(200, max_messages - len(raw_items))
            try:
                resp = client.conversations_replies(
                    channel=ref.channel_id,
                    ts=ref.thread_ts,
                    limit=limit,
                    cursor=cursor,
                    inclusive=True,
                )
            except SlackApiError as exc:
                err = exc.response.get("error", "unknown_error")
                if err == "not_in_channel" and not _auto_joined:
                    self._ensure_in_channel(client, ref.channel_id)
                    _auto_joined = True
                    continue  # retry with the joined channel
                if err == "not_in_channel":
                    raise ValueError(
                        f"Bot is not in channel '{ref.channel_id}'. "
                        "Invite the bot to the channel first."
                    ) from exc
                if err == "channel_not_found":
                    raise ValueError("Channel not found. Check the Slack thread URL.") from exc
                if err == "thread_not_found":
                    raise ValueError("Thread not found. Check the Slack thread URL.") from exc
                if err == "missing_scope":
                    raise RuntimeError(
                        "Missing Slack API scopes. Ensure the bot token has "
                        "'channels:history', 'groups:history', and 'channels:join' scopes."
                    ) from exc
                if err == "invalid_auth":
                    raise RuntimeError(
                        "Invalid Slack token. Check SLACK_BOT_TOKEN in backend/.env and recreate the backend container."
                    ) from exc
                raise RuntimeError(f"Slack API error: {err}") from exc

            for raw in resp.get("messages", []):
                if not include_bots and (raw.get("subtype") == "bot_message" or raw.get("bot_id")):
                    continue
                raw_items.append(raw)
                if len(raw_items) >= max_messages:
                    break

            cursor = (resp.get("response_metadata") or {}).get("next_cursor") or None
            if not cursor:
                break

        # ── Phase 2: Batch-resolve unique user IDs (parallel) ─────────────────
        unique_users = {
            raw.get("user") or raw.get("bot_id")
            for raw in raw_items
            if raw.get("user") or raw.get("bot_id")
        }
        users_to_resolve = [
            uid for uid in unique_users
            if uid and uid not in self._user_cache
        ]
        if users_to_resolve:
            with ThreadPoolExecutor(max_workers=min(8, len(users_to_resolve))) as pool:
                futures = {pool.submit(self._resolve_user, uid): uid for uid in users_to_resolve}
                for fut in as_completed(futures):
                    try:
                        fut.result()
                    except Exception:
                        pass  # _resolve_user already caches fallback

        # ── Phase 3: Collect all attachment file dicts for parallel download ──
        pending_files: List[Tuple[int, Dict]] = []
        for msg_idx, raw in enumerate(raw_items):
            for slack_file in raw.get("files", []) or []:
                pending_files.append((msg_idx, slack_file))

        # Download & process attachments concurrently
        processed_files: Dict[int, List[SlackThreadAttachment]] = {}
        all_attachments: List[SlackThreadAttachment] = []
        if pending_files:
            with ThreadPoolExecutor(max_workers=_MAX_ATTACHMENT_WORKERS) as pool:
                futures = {
                    pool.submit(self._process_attachment, sf): (mi, sf)
                    for mi, sf in pending_files
                }
                for fut in as_completed(futures):
                    mi, sf = futures[fut]
                    try:
                        item = fut.result()
                    except Exception as exc:
                        fname = sf.get("name", "unknown")
                        logger.warning("Attachment processing failed for %s: %s", fname, exc)
                        item = None
                    if item:
                        processed_files.setdefault(mi, []).append(item)
                        all_attachments.append(item)

        # ── Phase 4: Assemble SlackThreadMessage objects ──────────────────────
        messages: List[SlackThreadMessage] = []
        for msg_idx, raw in enumerate(raw_items):
            raw_text = raw.get("text") or ""
            clean_text, log_blocks = _extract_log_blocks(raw_text)
            message_attachments = processed_files.get(msg_idx, [])

            if not clean_text and not log_blocks and not message_attachments:
                continue

            ts = str(raw.get("ts", ""))
            try:
                dt = datetime.fromtimestamp(float(ts), tz=timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
            except Exception:
                dt = ts

            user_id = raw.get("user") or raw.get("bot_id")
            messages.append(
                SlackThreadMessage(
                    ts=ts,
                    datetime=dt,
                    user=self._resolve_user(user_id),
                    text=clean_text,
                    log_blocks=log_blocks,
                    attachments=message_attachments,
                )
            )

        return messages, all_attachments

    def _ollama_chat(self, chat_messages: List[Dict], model: str, max_tokens: int = 3500) -> str:
        # If LLMService is available, delegate to it for provider-aware routing
        if self._llm_service and hasattr(self._llm_service, "chat"):
            return self._llm_service.chat(
                messages=chat_messages,
                max_tokens=max_tokens,
                temperature=0.2,
                model_override=model if model != self._llm_service.model else None,
                module="slack_investigation",
            )

        # Fallback: direct Ollama API call (backward compatibility)
        payload = {
            "model": model,
            "messages": chat_messages,
            "stream": False,
            "options": {
                "temperature": 0.2,
                "num_ctx": settings.ollama_num_ctx,
            },
        }
        try:
            resp = requests.post(f"{self.ollama_host}/api/chat", json=payload, timeout=180)
            resp.raise_for_status()
            return (resp.json().get("message") or {}).get("content", "").strip()
        except requests.exceptions.ConnectionError as exc:
            raise RuntimeError(f"Ollama is not running at {self.ollama_host}. Start it with: ollama serve") from exc
        except Exception as exc:
            logger.error("Ollama chat error (model=%s): %s", model, exc)
            raise RuntimeError(f"Local LLM error: {exc}") from exc

    def _generate_summary(
        self,
        req: SlackThreadInvestigationRequest,
        messages: List[SlackThreadMessage],
        attachments: List[SlackThreadAttachment],
    ) -> Tuple[str, str]:
        # Determine model: use request override, active provider, or default
        model = req.model_override or ""

        # If LLMService is available and no override specified, use active provider
        if self._llm_service and hasattr(self._llm_service, "active_provider") and not model:
            active = self._llm_service.active_provider
            model = f"{active['type']}:{active['model']}"

        # If still no model, use the default text model
        if not model:
            model = self.text_model

        # For Ollama models (without provider prefix), validate installation
        is_openai = model.startswith("openai:")
        if not is_openai:
            plain_model = model.removeprefix("ollama:")
            installed = self._ollama_models()
            model_prefix = plain_model.split(":", 1)[0]
            if not any(model_prefix in name for name in installed):
                if any(self.text_model.split(":", 1)[0] in name for name in installed):
                    logger.warning(
                        "Model %s not installed, falling back to text model %s",
                        plain_model, self.text_model,
                    )
                    plain_model = self.text_model
                else:
                    raise RuntimeError(
                        f"Model '{plain_model}' is not installed in Ollama. "
                        f"Pull it with: ollama pull {plain_model}"
                    )
            model = plain_model

        system = (
            "You are a senior SRE producing a structured incident summary for a warehouse robotics team.\n"
            "You MUST summarize with respect to the ISSUE NAME and DESCRIPTION provided — every insight must be relevant to that specific issue.\n\n"
            "STRICT RULES:\n"
            "- Do NOT include any timestamps, dates, or times.\n"
            "- Do NOT include any URLs, rosbag links, or file paths.\n"
            "- Do NOT copy sentences from the thread; fully rephrase everything.\n"
            "- Do NOT mention user names or Slack handles; write impersonally.\n"
            "- Do NOT write long paragraphs; use ONLY concise bullet points (- item).\n"
            "- Focus every bullet on meaningful technical insight, not generic filler.\n\n"
            "Use EXACTLY these five markdown sections:\n\n"
            "## Issue Overview\n"
            "- (concise bullets describing the incident in context of the issue name/description)\n"
            "- (affected systems, components, or robots)\n\n"
            "## Key Observations\n"
            "- (specific technical findings from messages, logs, and attachments)\n"
            "- (quote exact error messages or log lines without timestamps)\n"
            "- (include service names, error codes, node names, metric values)\n\n"
            "## Root Cause Analysis\n"
            "- (contributing factors identified or suspected, tied to the issue)\n"
            "- (state uncertainty explicitly: 'Likely…', 'Unconfirmed…')\n\n"
            "## Actions Taken / Suggested Fixes\n"
            "- (actions performed during the incident)\n"
            "- (fixes applied or proposed)\n\n"
            "## Current Status / Risks\n"
            "- (resolution state: resolved, mitigated, or ongoing)\n"
            "- (remaining risks or recommended follow-ups)\n\n"
            "If information for a section is unavailable, write: - No information available in thread.\n"
        )

        # Strip timestamps and URLs from thread messages for the prompt
        thread_lines: List[str] = []
        for msg in messages:
            clean_msg = _URL_RE.sub("", msg.text).strip()
            line = f"{msg.user}: {clean_msg}"
            for idx, block in enumerate(msg.log_blocks, 1):
                clean_block = _URL_RE.sub("", block[:2000]).strip()
                line += f"\n  [Log block {idx}]\n{clean_block}"
            for attachment in msg.attachments:
                if attachment.extracted and attachment.filetype != "image":
                    line += f"\n  [Attachment: {attachment.filename} ({attachment.filetype})]\n{attachment.extracted[:3000]}"
                else:
                    line += f"\n  [Attachment: {attachment.filename} ({attachment.filetype})]"
            thread_lines.append(line)

        attachment_sections: List[str] = []
        for attachment in attachments:
            if attachment.filetype == "image":
                attachment_sections.append(f"=== IMAGE: {attachment.filename} ===\n[See visual content above]")
            else:
                attachment_sections.append(
                    f"=== {attachment.filetype.upper()}: {attachment.filename} ===\n{attachment.extracted}"
                )

        prompt = (
            f"ISSUE NAME / DESCRIPTION:\n{req.description}\n\n"
            f"SITE: {req.site_id or 'N/A'}\n\n"
            f"SLACK THREAD ({len(messages)} messages)\n"
            "--- MESSAGES ---\n\n"
            + "\n\n".join(thread_lines)
        )
        if attachment_sections:
            prompt += "\n\n" + ("-" * 60) + "\nATTACHMENTS\n" + ("-" * 60) + "\n\n"
            prompt += "\n\n".join(attachment_sections)
        if req.custom_prompt:
            prompt += f"\n\nSPECIAL FOCUS: {req.custom_prompt}"

        # Keep the prompt within the context window.
        # ~3.5 chars per token; reserve 800 tokens for output + system prompt.
        max_prompt_chars = max(500, (settings.ollama_num_ctx - 800) * 3)
        if len(prompt) > max_prompt_chars:
            prompt = prompt[:max_prompt_chars] + "\n\n[... thread truncated to fit context window ...]"

        chat: List[Dict] = [{"role": "system", "content": system}]
        user_message: Dict = {"role": "user", "content": prompt}
        chat.append(user_message)

        # Use lower max_tokens for OpenAI to improve response speed
        max_tok = 2000 if is_openai else 3500
        summary = self._ollama_chat(chat, model, max_tokens=max_tok)
        return summary, model

    def _infer_risk(self, summary: str) -> str:
        lowered = summary.lower()
        if any(token in lowered for token in ["sev1", "critical", "safety", "production down", "data loss"]):
            return "high"
        if any(token in lowered for token in ["degraded", "intermittent", "warning", "retry"]):
            return "medium"
        return "low"

    def investigate(self, req: SlackThreadInvestigationRequest) -> SlackThreadInvestigationResponse:
        # Determine if we're using OpenAI (skip Ollama-specific checks)
        using_openai = False
        if self._llm_service and hasattr(self._llm_service, "active_provider"):
            using_openai = self._llm_service.active_provider.get("type") == "openai"
        if req.model_override and req.model_override.startswith("openai:"):
            using_openai = True

        if not using_openai and not self._ollama_ping():
            raise RuntimeError(
                f"Ollama is not running at {self.ollama_host}. Run: ollama serve"
            )

        ref = parse_slack_thread_url(req.slack_thread_url)

        # Pre-warm Ollama model non-blocking (skip for OpenAI)
        if not using_openai:
            warm_model = req.model_override or self.text_model
            if warm_model.startswith("ollama:"):
                warm_model = warm_model.removeprefix("ollama:")
            def _warm():
                try:
                    requests.post(
                        f"{self.ollama_host}/api/generate",
                        json={"model": warm_model, "prompt": "", "keep_alive": "10m"},
                        timeout=5,
                    )
                except Exception:
                    pass
            # Fire-and-forget — model warms while we fetch messages
            import threading
            threading.Thread(target=_warm, daemon=True).start()

        messages, attachments = self._fetch_thread_messages(ref, req.include_bots, req.max_messages)
        if not messages:
            raise RuntimeError("Thread is empty or inaccessible with current token/scopes.")

        summary, model_used = self._generate_summary(req, messages, attachments)
        sections = _split_markdown_sections(summary)

        issue = _find_section(sections, "issue overview", "the issue", "issue", "problem", "incident")
        observations = _find_section(sections, "key observations", "observations", "important logs & errors", "important logs", "logs & errors")
        root_cause = _find_section(sections, "root cause analysis", "root cause", "cause")
        actions_taken = _find_section(sections, "actions taken / suggested fixes", "actions taken", "suggested fixes", "actions performed")
        status_risks = _find_section(sections, "current status / risks", "current status", "risks", "resolution & current status", "resolution", "status")

        # Build thread_summary from the key sections as bullet-point text
        summary_parts: List[str] = []
        if issue:
            summary_parts.append(f"**Issue Overview**\n{issue}")
        if observations:
            summary_parts.append(f"**Key Observations**\n{observations}")
        if root_cause:
            summary_parts.append(f"**Root Cause Analysis**\n{root_cause}")
        if actions_taken:
            summary_parts.append(f"**Actions Taken / Suggested Fixes**\n{actions_taken}")
        if status_risks:
            summary_parts.append(f"**Current Status / Risks**\n{status_risks}")

        thread_summary = "\n\n".join(summary_parts).strip() or summary[:2000]

        findings = _as_bullets(
            "\n".join(filter(None, [issue, observations, root_cause]))
        ) or [
            "Review raw analysis for detailed evidence extracted from messages and files."
        ]
        actions = _as_bullets(
            "\n".join(filter(None, [actions_taken, status_risks]))
        ) or [
            "No explicit action items detected; assign owners to follow up on unresolved findings."
        ]

        participants = sorted({msg.user for msg in messages if msg.user})

        usage = getattr(self._llm_service, "last_usage", {}) if self._llm_service else {}
        return SlackThreadInvestigationResponse(
            workspace=ref.workspace,
            channel_id=ref.channel_id,
            thread_ts=ref.thread_ts,
            message_count=len(messages),
            attachment_count=len(attachments),
            model_used=model_used,
            participants=participants,
            thread_summary=thread_summary,
            key_findings=findings,
            recommended_actions=actions,
            risk_level=self._infer_risk(summary),
            timeline=messages,
            attachments=attachments,
            raw_analysis=summary,
            actual_prompt_tokens     = usage.get("prompt_tokens", 0),
            actual_completion_tokens = usage.get("completion_tokens", 0),
            actual_total_tokens      = usage.get("total_tokens", 0),
            cost_usd                 = usage.get("cost_usd", 0.0),
        )
